"""
File Converter Module
Handles conversion between DOCX, PDF, PPTX, and Markdown formats.
"""

import sys
from pathlib import Path
from typing import Optional, List, Dict, Any
from dataclasses import dataclass
import re


@dataclass
class ConvertedDocument:
    """Result of document conversion."""
    text: str
    metadata: Dict[str, Any]
    images: List[str] = None
    tables: List[str] = None
    
    def __post_init__(self):
        if self.images is None:
            self.images = []
        if self.tables is None:
            self.tables = []


class FileConverter:
    """Converts between document formats."""
    
    def __init__(self):
        """Initialize converter."""
        pass
    
    def convert(self, input_path: str, output_format: str = "markdown") -> ConvertedDocument:
        """
        Convert a document to specified format.
        
        Args:
            input_path: Path to input file
            output_format: Target format (markdown, docx, pdf, text)
            
        Returns:
            ConvertedDocument object
        """
        path = Path(input_path)
        ext = path.suffix.lower()
        
        if ext == '.docx':
            return self._from_docx(input_path)
        elif ext == '.pdf':
            return self._from_pdf(input_path)
        elif ext == '.pptx':
            return self._from_pptx(input_path)
        elif ext == '.md':
            return self._from_markdown(input_path)
        else:
            raise ValueError(f"Unsupported format: {ext}")
    
    def _from_docx(self, input_path: str) -> ConvertedDocument:
        """Convert DOCX to markdown."""
        try:
            from docx import Document
            
            doc = Document(input_path)
            text_parts = []
            tables = []
            
            for para in doc.paragraphs:
                # Check for headings
                if para.style.name.startswith('Heading'):
                    level = para.style.name.split()[-1]
                    try:
                        level = int(level)
                        text_parts.append(f"{'#' * level} {para.text}")
                    except:
                        text_parts.append(f"## {para.text}")
                else:
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    table_text.append(row_text)
                tables.append("\n".join(table_text))
            
            # Combine
            full_text = "\n".join(text_parts)
            if tables:
                full_text += "\n\n## Tables\n\n" + "\n\n".join(tables)
            
            return ConvertedDocument(
                text=full_text,
                metadata={
                    "source": input_path,
                    "format": "docx",
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(tables)
                }
            )
            
        except ImportError:
            raise ImportError("python-docx not installed. Install with: pip install python-docx")
    
    def _from_pdf(self, input_path: str) -> ConvertedDocument:
        """Convert PDF to markdown."""
        try:
            import PyPDF2
            
            text_parts = []
            with open(input_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
            
            full_text = "\n\n".join(text_parts)
            
            return ConvertedDocument(
                text=full_text,
                metadata={
                    "source": input_path,
                    "format": "pdf",
                    "pages": len(text_parts)
                }
            )
            
        except ImportError:
            raise ImportError("PyPDF2 not installed. Install with: pip install PyPDF2")
    
    def _from_pptx(self, input_path: str) -> ConvertedDocument:
        """Convert PPTX to markdown."""
        try:
            from pptx import Presentation
            
            prs = Presentation(input_path)
            slides_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = [f"## Slide {i+1}\n"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text.append(shape.text)
                
                slides_text.append("\n".join(slide_text))
            
            full_text = "\n\n".join(slides_text)
            
            return ConvertedDocument(
                text=full_text,
                metadata={
                    "source": input_path,
                    "format": "pptx",
                    "slides": len(prs.slides)
                }
            )
            
        except ImportError:
            raise ImportError("python-pptx not installed. Install with: pip install python-pptx")
    
    def _from_markdown(self, input_path: str) -> ConvertedDocument:
        """Load markdown file."""
        with open(input_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return ConvertedDocument(
            text=text,
            metadata={
                "source": input_path,
                "format": "md",
                "words": len(text.split())
            }
        )
    
    def to_docx(self, markdown_text: str, output_path: str, title: str = "Manuscript") -> str:
        """
        Convert markdown to DOCX.
        
        Args:
            markdown_text: Markdown content
            output_path: Output file path
            title: Document title
            
        Returns:
            Path to output file
        """
        from docx import Document
        from docx.shared import Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.style import WD_STYLE_TYPE
        
        doc = Document()
        
        # Set up styles
        doc.styles['Normal'].font.name = 'Times New Roman'
        doc.styles['Normal'].font.size = Pt(12)
        
        # Title
        title_para = doc.add_heading(title, 0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Parse markdown and add content
        lines = markdown_text.split('\n')
        in_table = False
        table_rows = []
        
        for line in lines:
            # Skip empty lines
            if not line.strip():
                continue
            
            # Headings
            if line.startswith('#'):
                level = len(line.split()[0])
                text = line.lstrip('#').strip()
                if level == 1:
                    doc.add_heading(text, 0)
                elif level == 2:
                    doc.add_heading(text, 1)
                elif level == 3:
                    doc.add_heading(text, 2)
                else:
                    doc.add_heading(text, 3)
            
            # Tables
            elif line.startswith('|') and not in_table:
                in_table = True
                table_rows = [line]
            elif line.startswith('|') and in_table:
                table_rows.append(line)
            elif not line.startswith('|') and in_table:
                in_table = False
                # Create table
                if len(table_rows) >= 2:
                    cols = len(table_rows[0].split('|')) - 2
                    table = doc.add_table(rows=1, cols=max(cols, 1))
                    table.style = 'Table Grid'
                    # Header row
                    header_cells = table.rows[0].cells
                    headers = table_rows[0].strip('|').split('|')
                    for i, header in enumerate(headers[:cols]):
                        header_cells[i].text = header.strip()
                    # Data rows
                    for row_data in table_rows[1:]:
                        cells = row_data.strip('|').split('|')
                        row = table.add_row().cells
                        for i, cell_data in enumerate(cells[:cols]):
                            row[i].text = cell_data.strip()
                table_rows = []
            
            # Bullet points
            elif line.strip().startswith('- ') or line.strip().startswith('* '):
                doc.add_paragraph(line.strip()[2:], style='List Bullet')
            
            # Numbered lists
            elif re.match(r'^\d+\.\s', line):
                doc.add_paragraph(line, style='List Number')
            
            # Regular text
            else:
                # Clean up markdown formatting
                clean_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', line)  # Bold
                clean_text = re.sub(r'\*([^*]+)\*', r'\1', clean_text)  # Italic
                clean_text = re.sub(r'`([^`]+)`', r'\1', clean_text)   # Code
                doc.add_paragraph(clean_text)
        
        # Save
        doc.save(output_path)
        return output_path
    
    def extract_references(self, text: str) -> List[str]:
        """
        Extract references from markdown text.
        
        Args:
            text: Markdown text
            
        Returns:
            List of reference strings
        """
        refs = []
        lines = text.split('\n')
        in_refs = False
        
        for line in lines:
            if '## References' in line or '## REFERENCES' in line:
                in_refs = True
                continue
            if in_refs:
                if line.strip().startswith('[') or re.match(r'^\d+\.', line):
                    refs.append(line.strip())
                elif line and not line.startswith('#'):
                    refs.append(line.strip())
        
        return refs


def convert_document(input_path: str, output_format: str = "markdown") -> ConvertedDocument:
    """
    Convenience function to convert a document.
    
    Args:
        input_path: Path to input file
        output_format: Target format
        
    Returns:
        ConvertedDocument object
    """
    converter = FileConverter()
    return converter.convert(input_path, output_format)


def markdown_to_docx(markdown_text: str, output_path: str, title: str = "Manuscript") -> str:
    """
    Convert markdown to DOCX.
    
    Args:
        markdown_text: Markdown content
        output_path: Output file path
        title: Document title
        
    Returns:
        Path to output file
    """
    converter = FileConverter()
    return converter.to_docx(markdown_text, output_path, title)


if __name__ == "__main__":
    # Test
    import sys
    
    if len(sys.argv) < 3:
        print("Usage: python file_converter.py <input> <output> [--format md|docx]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    fmt = sys.argv[3] if len(sys.argv) > 3 else "markdown"
    
    converter = FileConverter()
    
    if fmt == "md":
        doc = converter.convert(input_file)
        with open(output_file, 'w') as f:
            f.write(doc.text)
        print(f"Saved markdown to: {output_file}")
    elif fmt == "docx":
        with open(input_file, 'r') as f:
            md_text = f.read()
        converter.to_docx(md_text, output_file)
        print(f"Saved DOCX to: {output_file}")
